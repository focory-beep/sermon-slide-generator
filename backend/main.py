from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from typing import List, Optional, Literal
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import tempfile
from datetime import datetime
from .bible_parser import parse_bible_reference
from .bible_loader import get_bible_loader
from .hymn_loader import get_hymn_loader
from .startup import initialize as startup_initialize

app = FastAPI(title="Sermon Slide Generator v3")

# 서버 시작 시 Reference 데이터 다운로드
@app.on_event("startup")
async def startup_event():
    """서버 시작 시 초기화 작업"""
    startup_initialize()

# 정적 파일 서빙 (프론트엔드)
frontend_path = os.path.join(os.path.dirname(__file__), "..", "frontend")
if os.path.exists(frontend_path):
    app.mount("/static", StaticFiles(directory=frontend_path), name="static")

# CORS 설정
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class ScriptureVerse(BaseModel):
    """성경 구절 모델"""
    reference: str  # 예: "마5:13" or "요한복음 3:16"
    text: str
    translation: str = "개역개정"
    reference_position: Literal["top-left", "top-right", "bottom-left", "bottom-right"] = "top-left"

class WorshipOrder(BaseModel):
    """예배 순서 항목 - 한 순서가 한 슬라이드"""
    title: str
    detail: Optional[str] = None

class SlideConfig(BaseModel):
    """슬라이드 디자인 설정 - 교회 예시 스타일"""
    font_name: str = "Noto Sans KR"  # Google Sans KR 스타일의 폰트
    font_size: int = 32  # 본문 폰트 크기 증가
    title_font_size: int = 44
    background_color: str = "#FFFFFF"
    text_color: str = "#000000"  # 더 진한 검정색
    title_color: str = "#000000"
    max_lines_per_slide: int = 10  # 슬라이드당 최대 줄 수 증가
    chars_per_line: int = 50  # 줄당 최대 글자 수 증가

class HymnRequest(BaseModel):
    """찬송가 요청 모델"""
    hymn_number: int

class PresentationRequest(BaseModel):
    """PPT 생성 요청"""
    title: str = "주일 예배"
    date: str = datetime.now().strftime("%Y년 %m월 %d일")
    worship_orders: List[WorshipOrder] = []
    scriptures: List[ScriptureVerse] = []
    hymns: List[HymnRequest] = []  # 찬송가 목록
    config: SlideConfig = SlideConfig()
    auto_parse_references: bool = True  # 자동 성경 참조 파싱

def hex_to_rgb(hex_color: str) -> RGBColor:
    """Hex 색상을 RGB로 변환"""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )

def add_title_slide(prs: Presentation, title: str, subtitle: str, config: SlideConfig):
    """제목 슬라이드 추가"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃

    # 배경색 설정
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(config.background_color)

    # 제목 텍스트박스 (중앙)
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(config.title_font_size + 12)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = config.font_name
    title_frame.paragraphs[0].font.color.rgb = hex_to_rgb(config.title_color)

    # 부제목 텍스트박스
    subtitle_box = slide.shapes.add_textbox(left, top + Inches(1.8), width, Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_frame.paragraphs[0].font.size = Pt(config.font_size)
    subtitle_frame.paragraphs[0].font.name = config.font_name
    subtitle_frame.paragraphs[0].font.color.rgb = hex_to_rgb(config.text_color)

def add_single_worship_order_slide(prs: Presentation, order: WorshipOrder, config: SlideConfig):
    """
    예배 순서 슬라이드 추가 - 한 순서가 한 슬라이드
    교회 예시와 비슷한 디자인: 중앙에 제목, 부제목
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경색
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(config.background_color)

    # 제목 (예배 순서 항목) - 중앙 배치
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = order.title
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(config.title_font_size)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = config.font_name
    title_frame.paragraphs[0].font.color.rgb = hex_to_rgb(config.title_color)

    # 상세 정보 (있는 경우)
    if order.detail:
        detail_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(7), Inches(1))
        detail_frame = detail_box.text_frame
        detail_frame.text = order.detail
        detail_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        detail_frame.paragraphs[0].font.size = Pt(config.font_size - 4)
        detail_frame.paragraphs[0].font.name = config.font_name
        detail_frame.paragraphs[0].font.color.rgb = hex_to_rgb(config.text_color)

def add_hymn_slides(prs: Presentation, hymn_number: int, config: SlideConfig):
    """
    찬송가 슬라이드 추가
    - 제목 슬라이드 (찬송가 번호 + 제목)
    - 각 절마다 별도 슬라이드
    - 후렴이 있으면 후렴 슬라이드 추가
    """
    loader = get_hymn_loader()
    hymn_data = loader.load_hymn(hymn_number)

    if not hymn_data:
        # 찬송가 데이터가 없으면 기본 슬라이드
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(config.background_color)

        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = f"찬송가 {hymn_number}장"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(config.title_font_size)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.name = config.font_name
        title_frame.paragraphs[0].font.color.rgb = hex_to_rgb(config.title_color)
        return

    # 1. 제목 슬라이드
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(config.background_color)

    # 찬송가 번호
    num_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(0.8))
    num_frame = num_box.text_frame
    num_frame.text = f"찬송가 {hymn_number}장"
    num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    num_frame.paragraphs[0].font.size = Pt(config.font_size)
    num_frame.paragraphs[0].font.name = config.font_name
    num_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = hymn_data['title']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(config.title_font_size)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = config.font_name
    title_frame.paragraphs[0].font.color.rgb = hex_to_rgb(config.title_color)

    # 2. 각 절 슬라이드
    for idx, verse in enumerate(hymn_data['verses'], 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(config.background_color)

        # 절 번호 (작게 상단)
        verse_num_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.5))
        verse_num_frame = verse_num_box.text_frame
        verse_num_frame.text = f"{idx}절"
        verse_num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        verse_num_frame.paragraphs[0].font.size = Pt(20)
        verse_num_frame.paragraphs[0].font.name = config.font_name
        verse_num_frame.paragraphs[0].font.color.rgb = RGBColor(120, 120, 120)

        # 가사
        lyrics_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4.5))
        lyrics_frame = lyrics_box.text_frame
        lyrics_frame.text = verse
        lyrics_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        lyrics_frame.paragraphs[0].font.size = Pt(config.font_size - 4)
        lyrics_frame.paragraphs[0].font.name = config.font_name
        lyrics_frame.paragraphs[0].font.color.rgb = hex_to_rgb(config.text_color)
        lyrics_frame.paragraphs[0].line_spacing = 1.8

    # 3. 후렴 슬라이드 (있는 경우)
    if hymn_data.get('chorus'):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(config.background_color)

        # 후렴 제목
        chorus_title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.5))
        chorus_title_frame = chorus_title_box.text_frame
        chorus_title_frame.text = "후렴"
        chorus_title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        chorus_title_frame.paragraphs[0].font.size = Pt(24)
        chorus_title_frame.paragraphs[0].font.bold = True
        chorus_title_frame.paragraphs[0].font.name = config.font_name
        chorus_title_frame.paragraphs[0].font.color.rgb = hex_to_rgb(config.title_color)

        # 후렴 가사
        chorus_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4.5))
        chorus_frame = chorus_box.text_frame
        chorus_frame.text = hymn_data['chorus']
        chorus_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        chorus_frame.paragraphs[0].font.size = Pt(config.font_size - 4)
        chorus_frame.paragraphs[0].font.name = config.font_name
        chorus_frame.paragraphs[0].font.color.rgb = hex_to_rgb(config.text_color)
        chorus_frame.paragraphs[0].line_spacing = 1.8

def split_scripture_text(text: str, max_lines: int, chars_per_line: int) -> List[str]:
    """
    성경 본문을 슬라이드에 맞게 분할
    줄 수와 줄당 글자 수를 고려하여 자연스럽게 분할
    """
    # 문장 단위로 분리 (마침표, 느낌표, 물음표 기준)
    import re
    sentences = re.split(r'([.!?]\s+|\d+\s+)', text)
    sentences = [''.join(sentences[i:i+2]) for i in range(0, len(sentences), 2)]

    chunks = []
    current_chunk = ""
    current_lines = 0

    for sentence in sentences:
        # 현재 문장이 몇 줄을 차지할지 추정
        sentence_lines = (len(sentence) + chars_per_line - 1) // chars_per_line

        if current_lines + sentence_lines <= max_lines:
            current_chunk += sentence
            current_lines += sentence_lines
        else:
            if current_chunk:
                chunks.append(current_chunk.strip())
            current_chunk = sentence
            current_lines = sentence_lines

    if current_chunk:
        chunks.append(current_chunk.strip())

    return chunks if chunks else [text]

def get_reference_position(position: str, slide_width=Inches(10), slide_height=Inches(7.5)):
    """
    레퍼런스 위치를 반환
    position: "top-left", "top-right", "bottom-left", "bottom-right"
    """
    margin = Inches(0.5)
    ref_width = Inches(3)
    ref_height = Inches(0.6)

    positions = {
        "top-left": (margin, margin),
        "top-right": (slide_width - ref_width - margin, margin),
        "bottom-left": (margin, slide_height - ref_height - margin),
        "bottom-right": (slide_width - ref_width - margin, slide_height - ref_height - margin)
    }

    return positions.get(position, positions["top-left"])

def add_scripture_slides(prs: Presentation, scripture: ScriptureVerse, config: SlideConfig, auto_parse: bool = True):
    """
    성경 구절 슬라이드 추가 - 교회 예시 스타일
    - 레퍼런스는 축약형으로 간결하게 (예: [마5:13])
    - 본문은 왼쪽 정렬로 깔끔하게 배치
    - 긴 본문은 자동으로 다음 슬라이드로 분할
    """
    # 성경 참조 파싱 (자동 매핑)
    reference_display = scripture.reference
    if auto_parse:
        try:
            lang = "korean"
            if scripture.translation.upper() in ["NIV", "ESV", "KJV", "NKJV", "NASB"]:
                lang = "english"
            elif "Luther" in scripture.translation or "Elberfelder" in scripture.translation:
                lang = "german"

            parsed = parse_bible_reference(scripture.reference, lang)
            # 교회 예시 스타일: 축약형 사용 (마5:13)
            book_abbrev = parsed.get("book_abbrev", scripture.reference)
            chapter = parsed.get("chapter")
            verses = parsed.get("verses")

            if book_abbrev and chapter:
                reference_display = f"{book_abbrev}{chapter}"
                if verses:
                    reference_display += f":{verses}"
            else:
                reference_display = scripture.reference
        except:
            reference_display = scripture.reference

    # 본문을 슬라이드에 맞게 분할
    text_chunks = split_scripture_text(
        scripture.text,
        config.max_lines_per_slide,
        config.chars_per_line
    )

    for i, chunk in enumerate(text_chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 배경색
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(config.background_color)

        # 성경 참조 위치 계산
        ref_left, ref_top = get_reference_position(scripture.reference_position)

        # 성경 참조 텍스트박스 - 교회 예시 스타일: 간결하게
        reference_text = f"[{reference_display}]"
        # 페이지 번호는 표시하지 않음 (교회 예시와 동일)

        ref_box = slide.shapes.add_textbox(ref_left, ref_top, Inches(4), Inches(0.8))
        ref_frame = ref_box.text_frame
        ref_frame.text = reference_text

        # 레퍼런스 정렬 (위치에 따라)
        if "right" in scripture.reference_position:
            ref_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        else:
            ref_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        ref_frame.paragraphs[0].font.size = Pt(24)  # 조금 더 크게
        ref_frame.paragraphs[0].font.bold = True  # 볼드체
        ref_frame.paragraphs[0].font.name = config.font_name
        ref_frame.paragraphs[0].font.color.rgb = hex_to_rgb(config.text_color)

        # 본문 텍스트박스 - 교회 예시 스타일: 왼쪽 정렬
        content_left = Inches(0.8)
        content_top = Inches(1.8)
        content_width = Inches(8.4)
        content_height = Inches(5)

        content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.text = chunk

        # 본문 스타일 - 교회 예시: 왼쪽 정렬
        content_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        content_frame.paragraphs[0].font.size = Pt(config.font_size)
        content_frame.paragraphs[0].font.name = config.font_name
        content_frame.paragraphs[0].font.color.rgb = hex_to_rgb(config.text_color)
        content_frame.paragraphs[0].line_spacing = 1.5

@app.post("/generate-presentation")
async def generate_presentation(request: PresentationRequest):
    """PPT 생성 API - 성경 본문 중심"""
    try:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # 1. 제목 슬라이드
        add_title_slide(prs, request.title, request.date, request.config)

        # 2. 예배 순서 슬라이드들 (한 순서 = 한 슬라이드)
        for order in request.worship_orders:
            add_single_worship_order_slide(prs, order, request.config)

        # 3. 성경 본문 슬라이드들 (자동 분할)
        for scripture in request.scriptures:
            add_scripture_slides(prs, scripture, request.config, request.auto_parse_references)

        # 4. 찬송가 슬라이드들
        for hymn in request.hymns:
            add_hymn_slides(prs, hymn.hymn_number, request.config)

        # 임시 파일로 저장
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        prs.save(temp_file.name)
        temp_file.close()

        return FileResponse(
            temp_file.name,
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            filename=f"{request.title}_{request.date}.pptx"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/", response_class=HTMLResponse)
async def root():
    """메인 페이지 - 웹 인터페이스"""
    try:
        html_path = os.path.join(os.path.dirname(__file__), "..", "frontend", "index_v3.html")
        if os.path.exists(html_path):
            with open(html_path, "r", encoding="utf-8") as f:
                html_content = f.read()
            html_content = html_content.replace("const API_URL = 'http://localhost:8000';", "const API_URL = window.location.origin;")
            return HTMLResponse(content=html_content)
        else:
            # fallback to v2
            html_path = os.path.join(os.path.dirname(__file__), "..", "frontend", "index_v2.html")
            if os.path.exists(html_path):
                with open(html_path, "r", encoding="utf-8") as f:
                    html_content = f.read()
                html_content = html_content.replace("const API_URL = 'http://localhost:8000';", "const API_URL = window.location.origin;")
                return HTMLResponse(content=html_content)
            return HTMLResponse(content="<h1>Frontend not found</h1>")
    except Exception as e:
        return HTMLResponse(content=f"<h1>Error loading frontend: {str(e)}</h1>")

@app.get("/api/info")
async def api_info():
    """API 정보"""
    return {
        "message": "Sermon Slide Generator API v3",
        "version": "3.0.0",
        "features": [
            "성경 본문 중심 슬라이드 생성",
            "레퍼런스 위치 선택 가능 (4곳)",
            "자동 성경 구절 매핑 및 파싱",
            "긴 본문 자동 분할",
            "한 예배 순서 = 한 슬라이드",
            "다국어 지원 (한국어, 영어, 독일어)",
            "커스터마이징 가능한 디자인"
        ]
    }

@app.get("/debug-reference")
async def debug_reference():
    """Reference 폴더 구조 디버깅"""
    import glob
    backend_dir = os.path.dirname(__file__)
    reference_dir = os.path.join(backend_dir, "Reference")

    result = {
        "backend_dir": backend_dir,
        "reference_dir": reference_dir,
        "reference_exists": os.path.exists(reference_dir),
        "files_in_backend": [],
        "files_in_reference": [],
        "bible_folder": None,
        "sample_files": []
    }

    # backend 디렉토리 파일 확인
    if os.path.exists(backend_dir):
        result["files_in_backend"] = os.listdir(backend_dir)[:20]

    # Reference 디렉토리 확인
    if os.path.exists(reference_dir):
        result["files_in_reference"] = os.listdir(reference_dir)[:20]

        # 개역개정 폴더 찾기
        for item in os.listdir(reference_dir):
            if '개역개정' in item or 'Bible' in item:
                result["bible_folder"] = item
                bible_path = os.path.join(reference_dir, item)
                if os.path.exists(bible_path):
                    result["sample_files"] = os.listdir(bible_path)[:10]
                break

    return result

@app.get("/parse-bible-reference")
async def parse_reference(reference: str, language: str = "korean"):
    """
    성경 참조 파싱 API
    예: /parse-bible-reference?reference=마5:13&language=korean
    """
    try:
        result = parse_bible_reference(reference, language)
        return result
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

@app.post("/auto-parse-scripture")
async def auto_parse_scripture(data: dict):
    """
    성경 구절 자동 파싱 API
    입력: {"reference": "마5:13", "language": "korean"}
    출력: {"book": "마태복음", "chapter": 5, "verse": 13, "formatted": "마태복음 5:13"}
    """
    try:
        reference = data.get("reference", "")
        language = data.get("language", "korean")
        result = parse_bible_reference(reference, language)
        return result
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

@app.post("/fetch-scripture")
async def fetch_scripture(data: dict):
    """
    성경 본문 자동 로드 API - 레퍼런스만 입력하면 본문을 가져옴
    입력: {"reference": "출24:12-18"}
    출력: {"reference": "출24:12-18", "text": "여호와께서 모세에게 이르시되...", "success": true}
    """
    try:
        reference = data.get("reference", "")
        if not reference:
            raise HTTPException(status_code=400, detail="Reference is required")

        print(f"📖 성경 본문 요청: {reference}")

        # 성경 본문 로드
        loader = get_bible_loader()
        print(f"   Bible loader path: {loader.bible_path}")
        print(f"   Path exists: {os.path.exists(loader.bible_path) if loader.bible_path else False}")

        text = loader.load_scripture(reference)

        if text:
            print(f"   ✅ 본문 로드 성공 ({len(text)} 글자)")
            return {
                "success": True,
                "reference": reference,
                "text": text
            }
        else:
            print(f"   ❌ 본문 로드 실패")
            # 디버깅 정보 추가
            backend_dir = os.path.dirname(__file__)
            ref_dir = os.path.join(backend_dir, "Reference")
            debug_info = {
                "backend_dir": backend_dir,
                "reference_dir": ref_dir,
                "ref_exists": os.path.exists(ref_dir),
                "loader_path": loader.bible_path,
                "loader_path_exists": os.path.exists(loader.bible_path) if loader.bible_path else False
            }
            if os.path.exists(ref_dir):
                debug_info["ref_contents"] = os.listdir(ref_dir)[:10]

            return {
                "success": False,
                "reference": reference,
                "error": "해당 구절을 찾을 수 없습니다. Reference 데이터가 설치되어 있는지 확인하세요.",
                "debug": debug_info
            }
    except Exception as e:
        print(f"   ❌ 오류 발생: {e}")
        import traceback
        traceback.print_exc()
        return {
            "success": False,
            "reference": data.get("reference", ""),
            "error": f"{type(e).__name__}: {str(e)}"
        }

@app.post("/fetch-hymn")
async def fetch_hymn(data: dict):
    """
    찬송가 데이터 자동 로드 API - 찬송가 번호로 제목과 가사 가져옴
    입력: {"hymn_number": 1}
    출력: {"number": 1, "title": "만복의 근원 하나님", "verses": [...], "success": true}
    """
    try:
        hymn_number = data.get("hymn_number")
        if not hymn_number:
            raise HTTPException(status_code=400, detail="Hymn number is required")

        hymn_number = int(hymn_number)
        if hymn_number < 1 or hymn_number > 645:
            raise HTTPException(status_code=400, detail="Hymn number must be between 1 and 645")

        # 찬송가 데이터 로드
        loader = get_hymn_loader()
        hymn_data = loader.load_hymn(hymn_number)

        if hymn_data:
            return {
                "success": True,
                **hymn_data
            }
        else:
            return {
                "success": False,
                "hymn_number": hymn_number,
                "error": "해당 찬송가를 찾을 수 없습니다. Reference 데이터가 설치되어 있는지 확인하세요."
            }
    except ValueError:
        return {
            "success": False,
            "error": "찬송가 번호는 숫자여야 합니다."
        }
    except Exception as e:
        return {
            "success": False,
            "hymn_number": data.get("hymn_number"),
            "error": str(e)
        }

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
