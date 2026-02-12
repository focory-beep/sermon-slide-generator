"""
독립 실행형 PPT 생성 데모
서버 없이 직접 PPT를 생성합니다.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def hex_to_rgb(hex_color: str) -> RGBColor:
    """Hex 색상을 RGB로 변환"""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )

def create_demo_presentation():
    """데모 프레젠테이션 생성"""
    print("📖 예배 슬라이드 데모 생성 중...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 설정
    font_name = "맑은 고딕"
    title_size = 44
    content_size = 32
    bg_color = hex_to_rgb("#FFFFFF")
    title_color = hex_to_rgb("#1a365d")
    text_color = hex_to_rgb("#000000")

    # 1. 제목 슬라이드
    print("  ✓ 제목 슬라이드 생성...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "주일 예배"
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(title_size + 8)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = font_name
    title_frame.paragraphs[0].font.color.rgb = title_color

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "2026년 2월 16일"
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_frame.paragraphs[0].font.size = Pt(content_size)
    subtitle_frame.paragraphs[0].font.name = font_name
    subtitle_frame.paragraphs[0].font.color.rgb = text_color

    # 2. 예배 순서 슬라이드
    print("  ✓ 예배 순서 슬라이드 생성...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "예배 순서"
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(title_size)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = font_name
    title_frame.paragraphs[0].font.color.rgb = title_color

    orders = [
        "예배로의 부름",
        "찬송 - 542장 구주 예수 의지함이",
        "기도 - 담임목사",
        "성경봉독 - 요한복음 3:16-21",
        "설교 - 하나님이 세상을 이처럼 사랑하사",
        "헌금 - 봉헌기도",
        "축도"
    ]

    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, order in enumerate(orders):
        p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
        p.text = f"{i+1}. {order}"
        p.font.size = Pt(content_size - 4)
        p.font.name = font_name
        p.font.color.rgb = text_color
        p.space_before = Pt(12)
        p.line_spacing = 1.5

    # 3. 성경 본문 슬라이드 (한국어)
    print("  ✓ 성경 본문 슬라이드 생성 (한국어)...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "요한복음 3:16-17"
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(title_size - 8)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = font_name
    title_frame.paragraphs[0].font.color.rgb = title_color

    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.text = "하나님이 세상을 이처럼 사랑하사 독생자를 주셨으니 이는 그를 믿는 자마다 멸망하지 않고 영생을 얻게 하려 하심이라 하나님이 그 아들을 세상에 보내신 것은 세상을 심판하려 하심이 아니요 그로 말미암아 세상이 구원을 받게 하려 하심이라"
    content_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    content_frame.paragraphs[0].font.size = Pt(content_size)
    content_frame.paragraphs[0].font.name = font_name
    content_frame.paragraphs[0].font.color.rgb = text_color
    content_frame.paragraphs[0].line_spacing = 1.6

    translation_box = slide.shapes.add_textbox(Inches(8), Inches(6.5), Inches(1.5), Inches(0.5))
    translation_frame = translation_box.text_frame
    translation_frame.text = "개역개정"
    translation_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    translation_frame.paragraphs[0].font.size = Pt(14)
    translation_frame.paragraphs[0].font.italic = True
    translation_frame.paragraphs[0].font.name = font_name
    translation_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)

    # 4. 성경 본문 슬라이드 (영어)
    print("  ✓ 성경 본문 슬라이드 생성 (English)...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "John 3:16-17"
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(title_size - 8)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = "Arial"
    title_frame.paragraphs[0].font.color.rgb = title_color

    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.text = "For God so loved the world that he gave his one and only Son, that whoever believes in him shall not perish but have eternal life. For God did not send his Son into the world to condemn the world, but to save the world through him."
    content_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    content_frame.paragraphs[0].font.size = Pt(content_size)
    content_frame.paragraphs[0].font.name = "Arial"
    content_frame.paragraphs[0].font.color.rgb = text_color
    content_frame.paragraphs[0].line_spacing = 1.6

    translation_box = slide.shapes.add_textbox(Inches(8), Inches(6.5), Inches(1.5), Inches(0.5))
    translation_frame = translation_box.text_frame
    translation_frame.text = "NIV"
    translation_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    translation_frame.paragraphs[0].font.size = Pt(14)
    translation_frame.paragraphs[0].font.italic = True
    translation_frame.paragraphs[0].font.name = "Arial"
    translation_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)

    # 5. 성경 본문 슬라이드 (독일어)
    print("  ✓ 성경 본문 슬라이드 생성 (Deutsch)...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Johannes 3:16-17"
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(title_size - 8)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = "Arial"
    title_frame.paragraphs[0].font.color.rgb = title_color

    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.text = "Denn so sehr hat Gott die Welt geliebt, dass er seinen eingeborenen Sohn gab, damit jeder, der an ihn glaubt, nicht verlorengeht, sondern ewiges Leben hat. Denn Gott hat seinen Sohn nicht in die Welt gesandt, damit er die Welt richte, sondern damit die Welt durch ihn gerettet werde."
    content_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    content_frame.paragraphs[0].font.size = Pt(content_size)
    content_frame.paragraphs[0].font.name = "Arial"
    content_frame.paragraphs[0].font.color.rgb = text_color
    content_frame.paragraphs[0].line_spacing = 1.6

    translation_box = slide.shapes.add_textbox(Inches(8), Inches(6.5), Inches(1.5), Inches(0.5))
    translation_frame = translation_box.text_frame
    translation_frame.text = "Luther 2017"
    translation_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    translation_frame.paragraphs[0].font.size = Pt(14)
    translation_frame.paragraphs[0].font.italic = True
    translation_frame.paragraphs[0].font.name = "Arial"
    translation_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)

    # 저장
    filename = "sermon_demo_multilingual.pptx"
    prs.save(filename)

    print(f"\n✅ 완료! 파일이 생성되었습니다: {filename}")
    print(f"📊 총 슬라이드 수: {len(prs.slides)}")
    print("\n슬라이드 목록:")
    print("  1. 제목 슬라이드 (주일 예배)")
    print("  2. 예배 순서")
    print("  3. 요한복음 3:16-17 (개역개정)")
    print("  4. John 3:16-17 (NIV)")
    print("  5. Johannes 3:16-17 (Luther 2017)")

    return filename

if __name__ == "__main__":
    print("="*60)
    print("  📖 예배 슬라이드 생성기 - 데모")
    print("  🌍 한국어 • English • Deutsch")
    print("="*60)
    print()

    create_demo_presentation()

    print()
    print("="*60)
    print("PowerPoint 파일을 열어서 결과를 확인해보세요! 🎉")
    print("="*60)
